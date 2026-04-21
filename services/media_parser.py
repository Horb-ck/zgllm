# /home/zgllm/test_server/services/media_parser.py
# -*- coding: utf-8 -*-
"""
多媒体文件解析服务
支持图片、视频、PPT 的智能解析

🔧 核心特性：
   - 多 VLM 模型自动回退：当第一个视觉模型不可用时自动切换下一个
   - 多 LLM 模型自动回退：文本摘要模型同理
   - ★ 模型成功缓存：记住上次成功的模型，下次优先使用
   - ★ 解析耗时统计：每个阶段独立计时，控制台输出详细报告
   - 若所有模型均失败 → 返回 success=False，前端显示「失败」
   - 详细的错误日志，便于排查

依赖（可选，缺少时对应功能不可用）：
   pip install opencv-python   # 视频帧提取
   pip install python-pptx     # PPT 解析
   pip install Pillow           # 图片处理
"""

import base64
import io
import os
import re
import tempfile
import time
import traceback
from typing import Dict, List, Optional, Callable, Any

import requests

# ================== 可选依赖检测 ==================
try:
    import cv2
    HAS_CV2 = True
except ImportError:
    HAS_CV2 = False

try:
    from pptx import Presentation as PptxPresentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False


class MediaParser:
    """
    多媒体文件解析器（带模型成功缓存 + 耗时统计）

    用法:
        parser = MediaParser(vlm_models=[...], llm_models=[...])
        result = parser.parse(file_bytes, 'photo.jpg', 'user01')
    """

    IMAGE_EXTENSIONS = {'jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp'}
    VIDEO_EXTENSIONS = {'mp4', 'avi', 'mov', 'mkv', 'flv', 'wmv', 'webm'}
    PPT_EXTENSIONS   = {'pptx', 'ppt'}
    ALL_EXTENSIONS   = IMAGE_EXTENSIONS | VIDEO_EXTENSIONS | PPT_EXTENSIONS

    # ============================================================
    # 初始化
    # ============================================================

    def __init__(
        self,
        vlm_models: List[Dict] = None,
        llm_models: List[Dict] = None,
        whisper_config: Dict = None,
    ):
        self.vlm_models = vlm_models or []
        self.llm_models = llm_models or []
        self.whisper_config = whisper_config or {}

        if not self.vlm_models:
            raise ValueError(
                "MediaParser 至少需要配置一个视觉语言模型 (vlm_models)。"
            )

        # ★★★ 模型成功缓存：记住上次成功的模型索引，下次优先使用 ★★★
        self._last_success_vlm_idx = 0
        self._last_success_llm_idx = 0

        # 启动日志
        print(f"   🖼️  VLM 模型回退链 ({len(self.vlm_models)} 个):")
        for i, m in enumerate(self.vlm_models):
            print(f"      {i+1}. {m['name']}  →  {m['model']}")
        print(f"   📝 LLM 模型回退链 ({len(self.llm_models)} 个):")
        for i, m in enumerate(self.llm_models):
            print(f"      {i+1}. {m['name']}  →  {m['model']}")

        cap_parts = []
        cap_parts.append(f"图片: {'✅' if True else '—'}")
        cap_parts.append(f"视频: {'✅' if HAS_CV2 else '❌ (需 pip install opencv-python)'}")
        cap_parts.append(f"PPT: {'✅' if HAS_PPTX else '❌ (需 pip install python-pptx)'}")
        print(f"   📋 功能: {' | '.join(cap_parts)}")

    # ============================================================
    # 公共接口
    # ============================================================

    def get_media_type(self, filename: str) -> str:
        ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
        if ext in self.IMAGE_EXTENSIONS:
            return 'image'
        if ext in self.VIDEO_EXTENSIONS:
            return 'video'
        if ext in self.PPT_EXTENSIONS:
            return 'ppt'
        return ''

    def parse(
        self,
        file_content: bytes,
        filename: str,
        username: str,
        on_progress: Callable[[str, int], None] = None,
    ) -> Dict[str, Any]:
        media_type = self.get_media_type(filename)
        if not media_type:
            return {'success': False, 'error': f'不支持的文件类型: {filename}'}

        file_size_mb = len(file_content) / 1024 / 1024                    # ★

        print(f"\n{'='*60}")
        print(f"🎬 开始多媒体解析: {filename} (类型: {media_type})")
        print(f"👤 用户: {username}")
        print(f"📦 文件大小: {len(file_content)} bytes ({file_size_mb:.2f} MB)")  # ★
        print(f"{'='*60}")

        # ★ 记录解析开始时间
        parse_start = time.time()

        try:
            if media_type == 'image':
                result = self._parse_image(file_content, filename, on_progress)
            elif media_type == 'video':
                result = self._parse_video(file_content, filename, on_progress)
            elif media_type == 'ppt':
                result = self._parse_ppt(file_content, filename, on_progress)
            else:
                result = {'success': False, 'error': f'未实现的媒体类型: {media_type}'}
        except Exception as e:
            traceback.print_exc()
            result = {'success': False, 'error': f'解析异常: {str(e)}'}

        # ★ 计算并打印解析耗时
        parse_elapsed = time.time() - parse_start
        result['parse_time'] = round(parse_elapsed, 2)

        print(f"\n{'─'*60}")
        print(f"⏱️  多媒体解析耗时报告")
        print(f"   📄 文件名:  {filename}")
        print(f"   🔧 类型:    {media_type}")
        print(f"   📦 大小:    {file_size_mb:.2f} MB")
        print(f"   ⏱️  总耗时:  {parse_elapsed:.2f}s")
        if file_size_mb > 0:
            print(f"   📊 速度:    {file_size_mb / parse_elapsed:.2f} MB/s" if parse_elapsed > 0 else "   📊 速度:    瞬时")
        print(f"   {'✅ 成功' if result.get('success') else '❌ 失败'}"
              f"{' — 文本 ' + str(len(result.get('text', ''))) + ' 字符' if result.get('success') and result.get('text') else ''}")
        print(f"{'─'*60}")

        return result

    # ============================================================
    # 图片解析
    # ============================================================

    def _parse_image(self, file_content, filename, on_progress=None):
        t0 = time.time()                                                   # ★

        if on_progress:
            on_progress('analyzing_frames', 10)

        file_content = self._maybe_compress_image(file_content, filename)
        data_url = self._bytes_to_data_url(file_content, filename)

        t_compress = time.time()                                           # ★
        print(f"   ⏱️  图片压缩/编码: {t_compress - t0:.2f}s")

        messages = [{
            'role': 'user',
            'content': [
                {
                    'type': 'text',
                    'text': (
                        '请详细描述这张图片的内容。要求：\n'
                        '1. 如果包含文字/公式/代码，请完整提取；\n'
                        '2. 如果包含图表/表格，请描述数据和结构；\n'
                        '3. 如果包含示意图/流程图，请描述逻辑关系；\n'
                        '4. 用 Markdown 格式组织输出。'
                    ),
                },
                {
                    'type': 'image_url',
                    'image_url': {'url': data_url},
                },
            ],
        }]

        if on_progress:
            on_progress('analyzing_frames', 30)

        try:
            t_vlm_start = time.time()                                      # ★
            description = self._call_vlm(messages, timeout=120)
            t_vlm_end = time.time()                                        # ★
            print(f"   ⏱️  VLM 推理耗时: {t_vlm_end - t_vlm_start:.2f}s")
        except RuntimeError as e:
            print(f"❌ 图片解析失败（所有 VLM 模型不可用）: {e}")
            return {'success': False, 'error': str(e)}

        if not description or not description.strip():
            return {'success': False, 'error': '视觉模型返回了空白内容'}

        if on_progress:
            on_progress('summarizing', 90)

        text = f"## 图片内容: {filename}\n\n{description}"

        if on_progress:
            on_progress('done', 100)

        print(f"✅ 图片解析完成，提取文本长度: {len(text)}")
        return {
            'success': True,
            'text': text,
            'metadata': {
                'media_type': 'image',
                'filename': filename,
                'text_length': len(text),
            },
        }

    # ============================================================
    # 视频解析
    # ============================================================

    def _parse_video(self, file_content, filename, on_progress=None):
        if not HAS_CV2:
            return {
                'success': False,
                'error': '视频解析需要安装 OpenCV：pip install opencv-python',
            }

        t0 = time.time()                                                   # ★

        if on_progress:
            on_progress('extracting_frames', 5)

        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(
                suffix=os.path.splitext(filename)[-1] or '.mp4',
                delete=False,
            ) as tmp:
                tmp.write(file_content)
                tmp_path = tmp.name

            cap = cv2.VideoCapture(tmp_path)
            if not cap.isOpened():
                return {'success': False, 'error': '无法打开视频文件'}

            fps = cap.get(cv2.CAP_PROP_FPS) or 30
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            duration = total_frames / fps if fps > 0 else 0

            max_frames = 20
            interval_sec = max(5, duration / max_frames)
            interval_frames = max(1, int(interval_sec * fps))

            frames: List[Dict] = []
            frame_idx = 0

            while True:
                ret, frame = cap.read()
                if not ret:
                    break

                if frame_idx % interval_frames == 0 or frame_idx == 0:
                    h, w = frame.shape[:2]
                    max_dim = 1024
                    if max(h, w) > max_dim:
                        scale = max_dim / max(h, w)
                        frame = cv2.resize(frame, (int(w * scale), int(h * scale)))

                    ok, buf = cv2.imencode('.jpg', frame, [cv2.IMWRITE_JPEG_QUALITY, 75])
                    if ok:
                        b64 = base64.b64encode(buf).decode('utf-8')
                        frames.append({
                            'base64': b64,
                            'timestamp': frame_idx / fps if fps > 0 else 0,
                        })

                    if len(frames) >= max_frames:
                        break

                frame_idx += 1

            cap.release()

            t_extract = time.time()                                        # ★
            print(f"   ⏱️  帧提取耗时: {t_extract - t0:.2f}s ({len(frames)} 帧)")

            if on_progress:
                on_progress('extracting_frames', 20)

            if not frames:
                return {'success': False, 'error': '无法从视频中提取到任何帧'}

            print(f"   📸 提取了 {len(frames)} 帧 (时长 {duration:.0f}s)")

            descriptions: List[str] = []
            all_vlm_failed = False
            t_vlm_total = 0.0                                              # ★

            for i, fr in enumerate(frames):
                if on_progress:
                    pct = 20 + int(60 * (i + 1) / len(frames))
                    on_progress('analyzing_frames', pct)

                ts = fr['timestamp']
                ts_str = f"{int(ts // 60):02d}:{int(ts % 60):02d}"
                data_url = f"data:image/jpeg;base64,{fr['base64']}"

                messages = [{
                    'role': 'user',
                    'content': [
                        {
                            'type': 'text',
                            'text': (
                                f'这是视频 {ts_str} 处的截图。'
                                f'请描述画面内容，完整提取所有可见的文字、图表和重要信息。'
                            ),
                        },
                        {
                            'type': 'image_url',
                            'image_url': {'url': data_url},
                        },
                    ],
                }]

                try:
                    t_frame_start = time.time()                            # ★
                    desc = self._call_vlm(messages, timeout=90)
                    t_frame_end = time.time()                              # ★
                    t_vlm_total += (t_frame_end - t_frame_start)
                    print(f"   ⏱️  第 {i+1}/{len(frames)} 帧 VLM: {t_frame_end - t_frame_start:.2f}s")

                    if desc and desc.strip():
                        descriptions.append(f"### 📍 [{ts_str}]\n\n{desc}")
                    else:
                        descriptions.append(f"### 📍 [{ts_str}]\n\n（模型返回空白）")
                except RuntimeError as e:
                    print(f"   ❌ 第 {i+1} 帧分析失败（所有 VLM 不可用）")
                    all_vlm_failed = True
                    break

            print(f"   ⏱️  VLM 推理总耗时: {t_vlm_total:.2f}s (共 {len(descriptions)} 帧)")  # ★

            if all_vlm_failed and not descriptions:
                return {
                    'success': False,
                    'error': '所有视觉模型均不可用，无法分析视频帧',
                }

            if on_progress:
                on_progress('summarizing', 85)

            combined_frames = '\n\n'.join(descriptions)
            summary = ''
            if self.llm_models and len(descriptions) >= 2:
                try:
                    t_sum_start = time.time()                              # ★
                    summary_messages = [{
                        'role': 'user',
                        'content': (
                            '以下是一段视频各时间点的画面描述。\n'
                            '请生成一份结构化的内容摘要，提取核心知识点，'
                            '使用 Markdown 格式。\n\n'
                            + combined_frames
                        ),
                    }]
                    summary = self._call_llm(summary_messages, timeout=90)
                    t_sum_end = time.time()                                # ★
                    print(f"   ⏱️  LLM 摘要耗时: {t_sum_end - t_sum_start:.2f}s")
                except RuntimeError:
                    print("   ⚠️ LLM 摘要生成失败（不影响主流程）")

            parts = [
                f"## 视频内容: {filename}\n",
                f"**时长**: {int(duration // 60):02d}:{int(duration % 60):02d}  ",
                f"**分析帧数**: {len(descriptions)}/{len(frames)}\n",
            ]

            if summary:
                parts.append(f"\n### 📋 内容摘要\n\n{summary}\n")

            parts.append(f"\n### 🎞️ 逐帧详细分析\n\n{combined_frames}")

            text = '\n'.join(parts)

            if on_progress:
                on_progress('done', 100)

            print(f"✅ 视频解析完成，提取文本长度: {len(text)}")
            return {
                'success': True,
                'text': text,
                'metadata': {
                    'media_type': 'video',
                    'filename': filename,
                    'duration': round(duration, 1),
                    'frames_total': len(frames),
                    'frames_analyzed': len(descriptions),
                    'has_summary': bool(summary),
                    'text_length': len(text),
                },
            }

        finally:
            if tmp_path:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass

    # ============================================================
    # PPT 解析
    # ============================================================

    def _parse_ppt(self, file_content, filename, on_progress=None):
        if not HAS_PPTX:
            return {
                'success': False,
                'error': 'PPT 解析需要安装 python-pptx：pip install python-pptx',
            }

        t0 = time.time()                                                   # ★

        if on_progress:
            on_progress('parsing_slides', 5)

        try:
            prs = PptxPresentation(io.BytesIO(file_content))
        except Exception as e:
            return {'success': False, 'error': f'无法打开 PPT 文件: {e}'}

        total_slides = len(prs.slides)
        if total_slides == 0:
            return {'success': False, 'error': 'PPT 文件没有任何幻灯片'}

        print(f"   📊 PPT 共 {total_slides} 页")

        slides_content: List[str] = []
        vlm_available = True
        t_vlm_total = 0.0                                                  # ★

        for i, slide in enumerate(prs.slides):
            if on_progress:
                pct = 5 + int(85 * (i + 1) / total_slides)
                on_progress('parsing_slides', pct)

            slide_texts: List[str] = []
            slide_images: List[bytes] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)

                if vlm_available:
                    try:
                        if hasattr(shape, 'image') and shape.image:
                            img_bytes = shape.image.blob
                            if img_bytes and len(img_bytes) > 500:
                                slide_images.append(img_bytes)
                    except Exception:
                        pass

            image_descs: List[str] = []
            if slide_images and vlm_available:
                for img_bytes in slide_images[:2]:
                    compressed = self._maybe_compress_image_bytes(
                        img_bytes, max_dim=1024, quality=80)
                    b64 = base64.b64encode(compressed).decode('utf-8')
                    data_url = f"data:image/jpeg;base64,{b64}"

                    messages = [{
                        'role': 'user',
                        'content': [
                            {
                                'type': 'text',
                                'text': '这是幻灯片中的图片，请描述其内容，提取所有文字和关键数据。',
                            },
                            {
                                'type': 'image_url',
                                'image_url': {'url': data_url},
                            },
                        ],
                    }]

                    try:
                        t_vlm_s = time.time()                              # ★
                        desc = self._call_vlm(messages, timeout=60)
                        t_vlm_e = time.time()                              # ★
                        t_vlm_total += (t_vlm_e - t_vlm_s)
                        print(f"   ⏱️  第 {i+1} 页图片 VLM: {t_vlm_e - t_vlm_s:.2f}s")

                        if desc and desc.strip():
                            image_descs.append(desc)
                    except RuntimeError:
                        print(f"   ⚠️ 第 {i+1} 页图片分析失败，后续页跳过图片分析")
                        vlm_available = False
                        break

            page_parts = [f"### 第 {i + 1} 页\n"]
            if slide_texts:
                page_parts.append('\n'.join(slide_texts))
            if image_descs:
                page_parts.append('\n**图片内容：**\n' + '\n'.join(image_descs))

            if slide_texts or image_descs:
                slides_content.append('\n'.join(page_parts))

        t_total = time.time() - t0                                         # ★
        print(f"   ⏱️  PPT 文本提取: {t_total - t_vlm_total:.2f}s, "      # ★
              f"VLM 图片分析: {t_vlm_total:.2f}s")

        if not slides_content:
            return {'success': False, 'error': 'PPT 中未提取到任何有效内容'}

        if on_progress:
            on_progress('summarizing', 95)

        text = (
            f"## PPT 内容: {filename}\n\n"
            f"**总页数**: {total_slides}\n\n"
            + '\n\n'.join(slides_content)
        )

        if on_progress:
            on_progress('done', 100)

        print(f"✅ PPT 解析完成，提取文本长度: {len(text)}，"
              f"有效页数: {len(slides_content)}/{total_slides}")

        return {
            'success': True,
            'text': text,
            'metadata': {
                'media_type': 'ppt',
                'filename': filename,
                'total_slides': total_slides,
                'slides_with_content': len(slides_content),
                'text_length': len(text),
            },
        }

    # ============================================================
    # ★★★ 核心：多模型自动回退 + 成功缓存 ★★★
    # ============================================================

    def _call_vlm(self, messages: List[Dict], timeout: int = 120) -> str:
        return self._call_with_fallback(
            models=self.vlm_models,
            messages=messages,
            timeout=timeout,
            label='VLM',
        )

    def _call_llm(self, messages: List[Dict], timeout: int = 60) -> str:
        return self._call_with_fallback(
            models=self.llm_models,
            messages=messages,
            timeout=timeout,
            label='LLM',
        )

    def _call_with_fallback(
        self,
        models: List[Dict],
        messages: List[Dict],
        timeout: int,
        label: str,
    ) -> str:
        """
        按顺序尝试每个模型；成功立即返回，失败则切换下一个。
        ★ 自动记住上次成功的模型索引，下次优先使用。
        所有模型都失败则抛出 RuntimeError。
        """
        if not models:
            raise RuntimeError(f"没有配置任何 {label} 模型")

        # ★ 获取上次成功的模型索引
        cache_attr = f'_last_success_{label.lower()}_idx'
        last_success_idx = getattr(self, cache_attr, 0)

        # ★ 构建优先尝试顺序：上次成功的排第一，其余保持原顺序
        if 0 < last_success_idx < len(models):
            order = [last_success_idx] + [i for i in range(len(models)) if i != last_success_idx]
        else:
            order = list(range(len(models)))

        errors: List[str] = []

        for rank, orig_idx in enumerate(order):
            config = models[orig_idx]
            name = config.get('name', f'model-{orig_idx}')
            model_id = config.get('model', 'unknown')

            # ★ 如果不是原始第一个，标注"优先"
            priority_tag = ''
            if rank == 0 and orig_idx != 0:
                priority_tag = ' ⚡优先(上次成功)'

            try:
                print(f"   🔄 [{label}] 尝试模型 {rank+1}/{len(models)}: "
                      f"{name} ({model_id}){priority_tag}")

                content = self._single_api_call(config, messages, timeout)

                if content and content.strip():
                    print(f"   ✅ [{label}] {name} 调用成功 "
                          f"(返回 {len(content)} 字)")
                    # ★ 记住成功的模型索引
                    setattr(self, cache_attr, orig_idx)
                    return content
                else:
                    err = f"{name}: 返回空响应"
                    errors.append(err)
                    print(f"   ⚠️ [{label}] {err}")

            except requests.exceptions.Timeout:
                err = f"{name}: 请求超时 ({timeout}s)"
                errors.append(err)
                print(f"   ⚠️ [{label}] {err}")

            except requests.exceptions.ConnectionError as e:
                err = f"{name}: 连接失败 ({e})"
                errors.append(err)
                print(f"   ⚠️ [{label}] {err}")

            except Exception as e:
                err = f"{name}: {type(e).__name__}: {str(e)[:200]}"
                errors.append(err)
                print(f"   ⚠️ [{label}] {err}")

            if rank < len(order) - 1:
                time.sleep(0.5)

        error_detail = '\n'.join(f"  {i+1}. {e}" for i, e in enumerate(errors))
        raise RuntimeError(
            f"所有 {label} 模型均调用失败 ({len(errors)}/{len(models)}):\n"
            f"{error_detail}"
        )

    def _single_api_call(self, config, messages, timeout):
        url = config['api_url']
        api_key = config['api_key']
        model = config['model']

        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/json',
        }

        payload = {
            'model': model,
            'messages': messages,
            'max_tokens': 4096,
            'temperature': 0.3,
        }

        response = requests.post(url, headers=headers, json=payload, timeout=timeout)

        if response.status_code != 200:
            body = ''
            try:
                body = response.text[:500]
            except Exception:
                pass
            raise Exception(f"HTTP {response.status_code}: {body}")

        result = response.json()

        if 'error' in result:
            err_obj = result['error']
            if isinstance(err_obj, dict):
                msg = err_obj.get('message', str(err_obj))
            else:
                msg = str(err_obj)
            raise Exception(f"API error: {msg}")

        choices = result.get('choices', [])
        if not choices:
            raise Exception("API 返回空 choices 列表")

        content = choices[0].get('message', {}).get('content', '')
        content = self._clean_response(content)
        return content

    # ============================================================
    # 工具方法
    # ============================================================

    def _clean_response(self, text: str) -> str:
        if not text:
            return ''
        text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL)
        return text.strip()

    def _bytes_to_data_url(self, img_bytes, filename):
        ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else 'jpeg'
        mime_map = {
            'jpg': 'image/jpeg', 'jpeg': 'image/jpeg',
            'png': 'image/png', 'gif': 'image/gif',
            'bmp': 'image/bmp', 'webp': 'image/webp',
        }
        mime = mime_map.get(ext, 'image/jpeg')
        b64 = base64.b64encode(img_bytes).decode('utf-8')
        return f"data:{mime};base64,{b64}"

    def _maybe_compress_image(self, img_bytes, filename, max_dim=2048, quality=85):
        if not HAS_PIL:
            return img_bytes
        if len(img_bytes) < 1 * 1024 * 1024:
            return img_bytes
        try:
            img = Image.open(io.BytesIO(img_bytes))
            if img.mode not in ('RGB', 'L'):
                img = img.convert('RGB')
            if max(img.size) > max_dim:
                ratio = max_dim / max(img.size)
                new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
                img = img.resize(new_size, Image.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format='JPEG', quality=quality)
            compressed = buf.getvalue()
            print(f"   🗜️ 图片压缩: {len(img_bytes)} → {len(compressed)} bytes")
            return compressed
        except Exception as e:
            print(f"   ⚠️ 图片压缩失败: {e}")
            return img_bytes

    def _maybe_compress_image_bytes(self, img_bytes, max_dim=1024, quality=80):
        if not HAS_PIL:
            return img_bytes
        if len(img_bytes) < 200 * 1024:
            return img_bytes
        try:
            img = Image.open(io.BytesIO(img_bytes))
            if img.mode not in ('RGB', 'L'):
                img = img.convert('RGB')
            if max(img.size) > max_dim:
                ratio = max_dim / max(img.size)
                new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
                img = img.resize(new_size, Image.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format='JPEG', quality=quality)
            return buf.getvalue()
        except Exception:
            return img_bytes
